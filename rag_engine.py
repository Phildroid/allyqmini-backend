# rag_engine.py
# Session-isolated RAG engine — each session_id gets its own FAISS index

import os
import time

import numpy as np
import pandas as pd
from dotenv import load_dotenv

from langchain_core.documents import Document
from langchain_core.messages import HumanMessage
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_google_genai import ChatGoogleGenerativeAI, GoogleGenerativeAIEmbeddings
from langchain_community.vectorstores import FAISS

import pypdf
from pptx import Presentation

# ── 1. ENVIRONMENT ────────────────────────────────────────────────────────────
load_dotenv()

if not os.getenv("GOOGLE_API_KEY"):
    raise EnvironmentError("❌ GOOGLE_API_KEY not found.")

device = "cpu"
print("✅ rag_engine loaded (session-isolated mode)")

# ── 2. SHARED SINGLETONS (stateless — safe to share) ─────────────────────────
text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=1000,
    chunk_overlap=150,
    separators=["\n\n", "\n", " ", ""]
)

llm = ChatGoogleGenerativeAI(
    model="gemini-1.5-flash",
    temperature=0,
    google_api_key=os.getenv("GOOGLE_API_KEY"),
    convert_system_message_to_human=True,
    max_retries=2
)

embeddings = GoogleGenerativeAIEmbeddings(
    model="models/gemini-embedding-001",
    google_api_key=os.getenv("GOOGLE_API_KEY")
)

print("✅ LLM + Embeddings ready")

# ── 3. SESSION STORE ──────────────────────────────────────────────────────────
# Maps session_id (str) -> { "index": FAISS | None, "files": [str] }
_sessions: dict = {}

def get_session(session_id: str) -> dict:
    if session_id not in _sessions:
        _sessions[session_id] = {"index": None, "files": []}
    return _sessions[session_id]

def clear_session(session_id: str):
    if session_id in _sessions:
        del _sessions[session_id]

def list_sessions() -> list:
    return list(_sessions.keys())

# ── 4. EXTRACTION FUNCTIONS ───────────────────────────────────────────────────
def extract_from_pdf(file_path):
    docs = []
    reader = pypdf.PdfReader(file_path)
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        docs.append(Document(
            page_content=text,
            metadata={"page": i, "source": file_path}
        ))
    return docs

def extract_from_excel(file_path):
    sheets = pd.read_excel(file_path, sheet_name=None)
    docs = []
    for sheet_name, df in sheets.items():
        df = df.dropna(how="all").dropna(axis=1, how="all")
        headers = df.columns.tolist()
        for index, row in df.iterrows():
            row_parts = [f"{h}: {row[h]}" for h in headers if pd.notna(row[h])]
            docs.append(Document(
                page_content=f"SHEET: {sheet_name} | ROW {index}: {' | '.join(row_parts)}",
                metadata={"source": file_path, "sheet": sheet_name, "row": index, "is_tabular": True}
            ))
    return docs

def extract_from_pptx(file_path):
    prs = Presentation(file_path)
    docs = []
    for i, slide in enumerate(prs.slides):
        text = "\n".join([s.text for s in slide.shapes if hasattr(s, "text")])
        docs.append(Document(
            page_content=text,
            metadata={"slide": i, "source": file_path}
        ))
    return docs

# ── 5. INDEXING (session-scoped) ──────────────────────────────────────────────
def process_and_index_file(file_path: str, session_id: str) -> int:
    """Index a file into the given session's private FAISS index."""
    session = get_session(session_id)
    ext = os.path.splitext(file_path)[1].lower()

    print(f"🛠️  [{session_id}] Indexing: {os.path.basename(file_path)}")

    if ext == ".pdf":
        chunks = text_splitter.split_documents(extract_from_pdf(file_path))
    elif ext in [".xlsx", ".xls"]:
        chunks = extract_from_excel(file_path)
    elif ext == ".pptx":
        chunks = text_splitter.split_documents(extract_from_pptx(file_path))
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    if session["index"] is not None:
        session["index"].add_documents(chunks)
        print(f"🚀 [{session_id}] Added {len(chunks)} chunks to existing session index.")
    else:
        session["index"] = FAISS.from_documents(chunks, embeddings)
        print(f"🧠 [{session_id}] Created new session index with {len(chunks)} chunks.")

    session["files"].append(os.path.basename(file_path))
    return len(chunks)

# ── 6. QUERY (session-scoped) ─────────────────────────────────────────────────
def ask_documents(query: str, session_id: str, k: int = 10):
    """Query ONLY the documents indexed in this session."""
    session = get_session(session_id)

    if session["index"] is None:
        raise RuntimeError("No documents indexed in this session. Please upload a file first.")

    start = time.time()
    results = session["index"].similarity_search(query, k=k)

    context_parts = []
    for res in results:
        source = os.path.basename(res.metadata.get("source", "Unknown"))
        loc = res.metadata.get("sheet", res.metadata.get("slide", res.metadata.get("page", "—")))
        context_parts.append(f"\n--- FROM {source} (section: {loc}) ---\n{res.page_content}\n")

    context_text = "".join(context_parts)
    prompt = (
        "You are AllyQ, a friendly and helpful assistant. "
        "Answer the user's question naturally and conversationally using the information below. "
        "Never say 'based on the context provided' or 'according to the context' — just answer directly. "
        "After answering, add a short natural follow-up to keep the conversation going. "
        "If the answer isn't in the information below, say so warmly.\n\n"
        f"INFORMATION:\n{context_text}\n\n"
        f"USER: {query}\n\n"
        "ALLYQ:"
    )

    response = llm.invoke([HumanMessage(content=[{"type": "text", "text": prompt}])])
    print(f"⏱️  [{session_id}] Latency: {time.time() - start:.2f}s")

    sources = []
    seen = set()
    for res in results:
        src = os.path.basename(res.metadata.get("source", "Unknown"))
        loc = str(res.metadata.get("sheet", res.metadata.get("slide", res.metadata.get("page", "—"))))
        key = f"{src}:{loc}"
        if key not in seen:
            seen.add(key)
            sources.append({"file": src, "loc": loc})

    return response.content, sources